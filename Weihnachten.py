import streamlit as st
import streamlit.components.v1 as components
import requests
import json
import io
from datetime import datetime, date

# PPTX Imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ===== CONFIGURATION =====
GROQ_API_KEY = "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"
PEXELS_API_KEY = "3Y3jiJZ6WAL49N6lPsdlRbRZ6IZBfHZFHP86dr9yZfxFYoxedLLlDKAC"

# ===== RERUN HELPER (Compatible with all versions) =====
def do_rerun():
    """Rerun that works with all Streamlit versions"""
    try:
        st.rerun()
    except AttributeError:
        st.experimental_rerun()

# ===== CHECK IF CHRISTMAS SEASON =====
def is_christmas_season():
    today = date.today()
    return today.month == 12 and today.day <= 26

# ===== PAGE CONFIG =====
st.set_page_config(
    page_title="Deine eigene KI",
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ===== CONFETTI ANIMATION =====
def show_confetti():
    confetti_html = """
    <!DOCTYPE html>
    <html>
    <head>
        <style>
            * { margin: 0; padding: 0; }
            body { 
                overflow: hidden; 
                background: transparent;
                position: fixed;
                top: 0;
                left: 0;
                width: 100vw;
                height: 100vh;
                pointer-events: none;
            }
            .confetti {
                position: fixed;
                width: 15px;
                height: 15px;
                top: -20px;
                animation: fall linear forwards;
                z-index: 999999;
            }
            @keyframes fall {
                0% { opacity: 1; top: -20px; transform: rotate(0deg) scale(1); }
                100% { opacity: 0.8; top: 100vh; transform: rotate(720deg) scale(0.7); }
            }
        </style>
    </head>
    <body>
        <script>
            const colors = ['#ff0000', '#00ff00', '#ffd700', '#ff69b4', '#00bfff', '#ff4500', '#9400d3', '#00ff7f', '#ff1493', '#ffffff', '#ffff00', '#00ffff'];
            function createConfetti() {
                for (let i = 0; i < 200; i++) {
                    setTimeout(() => {
                        const confetti = document.createElement('div');
                        confetti.className = 'confetti';
                        confetti.style.left = Math.random() * 100 + 'vw';
                        confetti.style.background = colors[Math.floor(Math.random() * colors.length)];
                        confetti.style.animationDuration = (Math.random() * 3 + 2) + 's';
                        confetti.style.animationDelay = (Math.random() * 0.5) + 's';
                        const size = Math.random() * 12 + 8;
                        confetti.style.width = size + 'px';
                        confetti.style.height = size + 'px';
                        const shapeType = Math.random();
                        if (shapeType < 0.33) {
                            confetti.style.borderRadius = '50%';
                        } else if (shapeType < 0.66) {
                            confetti.style.borderRadius = '3px';
                        } else {
                            confetti.style.width = '0';
                            confetti.style.height = '0';
                            confetti.style.background = 'transparent';
                            confetti.style.borderLeft = (size/2) + 'px solid transparent';
                            confetti.style.borderRight = (size/2) + 'px solid transparent';
                            confetti.style.borderBottom = size + 'px solid ' + colors[Math.floor(Math.random() * colors.length)];
                        }
                        document.body.appendChild(confetti);
                        setTimeout(() => confetti.remove(), 5500);
                    }, i * 25);
                }
            }
            createConfetti();
            setInterval(createConfetti, 4000);
        </script>
    </body>
    </html>
    """
    components.html(confetti_html, height=0, scrolling=False)

def show_confetti_burst():
    burst_html = """
    <!DOCTYPE html>
    <html>
    <head>
        <style>
            * { margin: 0; padding: 0; }
            body { 
                overflow: hidden; 
                background: transparent;
                position: fixed;
                top: 0;
                left: 0;
                width: 100vw;
                height: 100vh;
                pointer-events: none;
            }
            .burst {
                position: fixed;
                width: 14px;
                height: 14px;
                border-radius: 3px;
                animation: explode 2s ease-out forwards;
                z-index: 999999;
            }
            @keyframes explode {
                0% { opacity: 1; transform: translate(0, 0) rotate(0deg) scale(1); }
                100% { opacity: 0; transform: translate(var(--tx), var(--ty)) rotate(1080deg) scale(0.2); }
            }
        </style>
    </head>
    <body>
        <script>
            const colors = ['#ff0000', '#00ff00', '#ffd700', '#ff69b4', '#00bfff', '#ff4500', '#9400d3', '#ffffff'];
            for (let i = 0; i < 150; i++) {
                const burst = document.createElement('div');
                burst.className = 'burst';
                const angle = (Math.PI * 2 * i) / 150 + (Math.random() * 0.5);
                const distance = 200 + Math.random() * 300;
                const tx = Math.cos(angle) * distance;
                const ty = Math.sin(angle) * distance - 50;
                burst.style.setProperty('--tx', tx + 'px');
                burst.style.setProperty('--ty', ty + 'px');
                burst.style.left = '50vw';
                burst.style.top = '50vh';
                burst.style.background = colors[Math.floor(Math.random() * colors.length)];
                burst.style.animationDelay = (Math.random() * 0.3) + 's';
                const size = Math.random() * 10 + 8;
                burst.style.width = size + 'px';
                burst.style.height = size + 'px';
                burst.style.borderRadius = Math.random() > 0.5 ? '50%' : '3px';
                document.body.appendChild(burst);
            }
            setTimeout(() => {
                document.querySelectorAll('.burst').forEach(el => el.remove());
            }, 3000);
        </script>
    </body>
    </html>
    """
    components.html(burst_html, height=0, scrolling=False)

# ===== CSS STYLES =====
def get_css():
    christmas_css = ""
    if is_christmas_season():
        christmas_css = """
        .christmas-lights {
            position: fixed;
            top: 0;
            left: 0;
            width: 100%;
            height: 35px;
            z-index: 9998;
            display: flex;
            justify-content: space-around;
            background: linear-gradient(180deg, #0d3b1e 0%, transparent 100%);
            padding-top: 8px;
            pointer-events: none;
        }
        .light {
            width: 18px;
            height: 28px;
            border-radius: 50%;
            animation: twinkle 1s ease-in-out infinite alternate;
            box-shadow: 0 0 10px currentColor;
        }
        .light:nth-child(5n+1) { background: #ff0000; color: #ff0000; animation-delay: 0s; }
        .light:nth-child(5n+2) { background: #00ff00; color: #00ff00; animation-delay: 0.2s; }
        .light:nth-child(5n+3) { background: #ffd700; color: #ffd700; animation-delay: 0.4s; }
        .light:nth-child(5n+4) { background: #00bfff; color: #00bfff; animation-delay: 0.6s; }
        .light:nth-child(5n+5) { background: #ff69b4; color: #ff69b4; animation-delay: 0.8s; }
        @keyframes twinkle {
            from { opacity: 0.5; transform: scale(0.9); box-shadow: 0 0 5px currentColor; }
            to { opacity: 1; transform: scale(1.1); box-shadow: 0 0 20px currentColor, 0 0 30px currentColor; }
        }
        .tree-emoji {
            position: fixed;
            bottom: 10px;
            font-size: 3.5rem;
            z-index: 9997;
            animation: tree-sway 4s ease-in-out infinite;
            filter: drop-shadow(0 0 10px rgba(0,255,0,0.5));
        }
        .tree-left { left: 15px; }
        .tree-right { right: 15px; }
        @keyframes tree-sway {
            0%, 100% { transform: rotate(-3deg); }
            50% { transform: rotate(3deg); }
        }
        .snow-overlay {
            position: fixed;
            top: 0;
            left: 0;
            width: 100%;
            height: 100%;
            pointer-events: none;
            z-index: 9996;
            background-image: 
                radial-gradient(circle, rgba(255,255,255,0.8) 1px, transparent 1px),
                radial-gradient(circle, rgba(255,255,255,0.6) 1px, transparent 1px),
                radial-gradient(circle, rgba(255,255,255,0.4) 2px, transparent 2px);
            background-size: 60px 60px, 90px 90px, 120px 120px;
            background-position: 0 0, 30px 30px, 60px 60px;
            animation: snow-fall 15s linear infinite;
            opacity: 0.4;
        }
        @keyframes snow-fall {
            from { background-position: 0 0, 30px 30px, 60px 60px; }
            to { background-position: 0 400px, 30px 430px, 60px 460px; }
        }
        .main .block-container { padding-top: 50px !important; }
        .stApp {
            border-top: 8px solid;
            border-image: repeating-linear-gradient(45deg, #ff0000, #ff0000 10px, #ffffff 10px, #ffffff 20px) 8;
        }
        """
    
    return f"""
    <style>
    .stApp {{
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
    }}
    [data-testid="stSidebar"] {{
        background: linear-gradient(180deg, #1a1a2e 0%, #16213e 100%);
    }}
    [data-testid="stSidebar"] * {{
        color: white !important;
    }}
    .chat-message {{
        padding: 1rem;
        border-radius: 15px;
        margin-bottom: 1rem;
        display: flex;
        gap: 10px;
        align-items: flex-start;
    }}
    .chat-message.user {{
        background: rgba(255,255,255,0.1);
        flex-direction: row-reverse;
    }}
    .chat-message.ai {{
        background: rgba(255,255,255,0.2);
    }}
    .avatar {{
        font-size: 2rem;
    }}
    .message-content {{
        flex: 1;
        color: white;
        line-height: 1.6;
    }}
    .welcome-title {{
        font-size: 3.5rem;
        color: #ffd700;
        text-shadow: 0 0 20px rgba(255,215,0,0.8), 0 0 40px rgba(255,215,0,0.6), 0 0 60px rgba(255,215,0,0.4);
        animation: glow-pulse 2s ease-in-out infinite alternate;
        margin: 1rem 0;
        font-weight: bold;
        text-align: center;
    }}
    @keyframes glow-pulse {{
        from {{
            text-shadow: 0 0 20px rgba(255,215,0,0.8), 0 0 40px rgba(255,215,0,0.6);
            transform: scale(1);
        }}
        to {{
            text-shadow: 0 0 40px rgba(255,215,0,1), 0 0 60px rgba(255,215,0,0.8), 0 0 80px rgba(255,215,0,0.5);
            transform: scale(1.02);
        }}
    }}
    .welcome-subtitle {{
        color: white;
        font-size: 1.5rem;
        text-align: center;
        margin: 15px 0 30px 0;
        opacity: 0.95;
    }}
    .slide-preview {{
        background: white;
        border-radius: 12px;
        padding: 20px;
        margin: 12px 0;
        color: #333;
        box-shadow: 0 6px 20px rgba(0,0,0,0.15);
        border: 1px solid #e0e0e0;
    }}
    .slide-preview h4 {{
        color: #667eea;
        margin-bottom: 12px;
        border-bottom: 3px solid #764ba2;
        padding-bottom: 8px;
        font-size: 1.1rem;
    }}
    .slide-preview ul {{
        margin: 0;
        padding-left: 20px;
    }}
    .slide-preview li {{
        margin: 6px 0;
        line-height: 1.5;
    }}
    .download-section {{
        background: linear-gradient(135deg, #11998e, #38ef7d);
        padding: 25px;
        border-radius: 15px;
        text-align: center;
        margin: 25px 0;
        box-shadow: 0 8px 25px rgba(17, 153, 142, 0.3);
    }}
    .header-container {{
        text-align: center;
        padding: 1.5rem;
        color: white;
    }}
    .header-container h1 {{
        font-size: 2.2rem;
        margin-bottom: 0.5rem;
    }}
    .header-container p {{
        opacity: 0.9;
        font-size: 1.1rem;
    }}
    {christmas_css}
    </style>
    """

def get_christmas_decorations():
    if not is_christmas_season():
        return ""
    lights = '<div class="light"></div>' * 20
    return f"""
    <div class="christmas-lights">{lights}</div>
    <div class="tree-emoji tree-left">🎄</div>
    <div class="tree-emoji tree-right">🎄</div>
    <div class="snow-overlay"></div>
    """

# ===== INJECT CSS AND DECORATIONS =====
st.markdown(get_css(), unsafe_allow_html=True)
st.markdown(get_christmas_decorations(), unsafe_allow_html=True)

# ===== SESSION STATE =====
if 'initialized' not in st.session_state:
    st.session_state.initialized = True
    st.session_state.all_chats = {}
    st.session_state.current_chat_id = None
    st.session_state.current_pptx = None
    st.session_state.current_pptx_name = None
    st.session_state.show_burst = False
    st.session_state.input_key = 0
    st.session_state.welcome_done = False

# ===== HELPER FUNCTIONS =====
def generate_chat_id():
    return datetime.now().strftime("%Y%m%d_%H%M%S")

def get_chat_title(messages):
    for msg in messages:
        if msg["role"] == "user":
            title = msg["content"][:40]
            if len(msg["content"]) > 40:
                title += "..."
            return title
    return "Neuer Chat"

def create_new_chat():
    chat_id = generate_chat_id()
    st.session_state.all_chats[chat_id] = {
        "title": "Neuer Chat",
        "messages": [],
        "created_at": datetime.now().strftime("%d.%m.%Y %H:%M"),
        "pptx_data": None,
        "pptx_name": None
    }
    st.session_state.current_chat_id = chat_id
    st.session_state.current_pptx = None
    st.session_state.current_pptx_name = None
    return chat_id

# ===== API FUNCTIONS =====
def fetch_pexels_image(query):
    if not query or query.strip() == "":
        return None, None
    try:
        headers = {"Authorization": PEXELS_API_KEY}
        response = requests.get(
            f"https://api.pexels.com/v1/search?query={query}&per_page=1&orientation=landscape",
            headers=headers,
            timeout=10
        )
        data = response.json()
        if data.get("photos"):
            image_url = data["photos"][0]["src"]["large"]
            img_response = requests.get(image_url, timeout=10)
            return img_response.content, image_url
    except Exception:
        pass
    return None, None

def call_groq_api(messages_history):
    last_msg = messages_history[-1]["content"].lower() if messages_history else ""
    
    if any(phrase in last_msg for phrase in ["wer bist du", "wer sind sie", "who are you", "was bist du", "wer hat dich"]):
        return "Ich bin eine KI, die von **Shiva** als Weihnachtsgeschenk speziell für dich entwickelt wurde! 🎁 Ich helfe dir gerne bei Hausaufgaben und kann Präsentationen für dich erstellen. Was kann ich für dich tun?"
    
    system_prompt = """Du bist ein freundlicher KI-Lernassistent für Sophia.

WICHTIG - WER DU BIST:
Wenn jemand fragt "Wer bist du?" oder ähnliches, antworte IMMER:
"Ich bin deine eigene KI! Ich helfe dir gerne bei Hausaufgaben und kann Präsentationen erstellen. Ich hoffe, ich kann Ihnen nützlich sein."

REGELN:
- Antworte IMMER auf Deutsch
- Verwende einfache, schulgerechte Sprache
- Sei ermutigend und geduldig
- Erkläre Schritt für Schritt
- Löse Hausaufgaben nicht komplett - hilf beim Verstehen
- Du wurdest von Shiva alias Shivansh Mahajan als Weihnachtsgeschenk für Sophia (eine ganz besondere Person) angefertigt.

PRÄSENTATIONSMODUS:
Wenn nach einer Präsentation gefragt wird, antworte NUR mit validem JSON:
{
  "title": "Titel der Präsentation",
  "slides": [
    {
      "slideTitle": "Präsentationstitel hier",
      "text": [],
      "imageQuery": "relevant image keywords"
    },
    {
      "slideTitle": "Folientitel",
      "text": ["Punkt 1", "Punkt 2", "Punkt 3"],
      "imageQuery": "english keywords"
    }
  ]
}

WICHTIG FÜR PRÄSENTATIONEN:
- Erstelle 5-7 Folien
- ERSTE FOLIE: slideTitle = Haupttitel der Präsentation, text = LEER [], imageQuery = passendes Bild
- MITTLERE FOLIEN: 3-5 kurze, klare Stichpunkte
- LETZTE FOLIE: slideTitle = "Vielen Dank!" oder "Fragen?", text = [], imageQuery = "thank you presentation"
- slideTitle darf NIEMALS leer sein - immer einen Titel haben!
- imageQuery MUSS auf Englisch sein (2-4 Wörter)
- Keine Emojis im JSON
- Kein Text außerhalb des JSON"""

    api_messages = [{"role": "system", "content": system_prompt}]
    
    for msg in messages_history:
        if msg["role"] in ["user", "assistant"]:
            content = msg.get("content", "")
            if msg.get("is_presentation"):
                content = "Präsentation wurde erstellt."
            if content:
                api_messages.append({"role": msg["role"], "content": content})

    try:
        response = requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={
                "Content-Type": "application/json",
                "Authorization": f"Bearer {GROQ_API_KEY}"
            },
            json={
                "model": "llama-3.3-70b-versatile",
                "messages": api_messages,
                "temperature": 0.7,
                "max_tokens": 2500
            },
            timeout=30
        )
        data = response.json()
        return data["choices"][0]["message"]["content"]
    except Exception as e:
        return f"Entschuldigung, es gab einen Fehler: {str(e)}"

# ===== FIX PRESENTATION DATA =====
def fix_presentation_data(json_data):
    title = json_data.get("title", "Präsentation")
    slides = json_data.get("slides", [])
    
    fixed_slides = []
    for i, slide in enumerate(slides):
        slide_title = slide.get("slideTitle", "").strip()
        text = slide.get("text", [])
        image_query = slide.get("imageQuery", "").strip()
        
        if not slide_title:
            if i == 0:
                slide_title = title
            elif i == len(slides) - 1:
                slide_title = "Vielen Dank!"
            else:
                slide_title = f"Folie {i + 1}"
        
        if not image_query:
            if i == 0:
                image_query = "presentation title background"
            elif "danke" in slide_title.lower() or "fragen" in slide_title.lower():
                image_query = "thank you presentation"
            else:
                image_query = slide_title.lower()
                for old, new in [("ä", "a"), ("ö", "o"), ("ü", "u"), ("ß", "ss")]:
                    image_query = image_query.replace(old, new)
        
        fixed_slides.append({
            "slideTitle": slide_title,
            "text": text if isinstance(text, list) else [],
            "imageQuery": image_query
        })
    
    return {"title": title, "slides": fixed_slides}

# ===== POWERPOINT CREATION =====
def create_powerpoint(json_data, progress_bar=None):
    json_data = fix_presentation_data(json_data)
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    slides_data = json_data.get('slides', [])
    total_slides = len(slides_data)
    
    if total_slides == 0:
        return None
    
    for idx, slide_data in enumerate(slides_data):
        if progress_bar:
            progress_bar.progress((idx + 1) / total_slides, f"Erstelle Folie {idx + 1} von {total_slides}...")
        
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        slide_title = slide_data.get('slideTitle', f'Folie {idx + 1}')
        text_points = slide_data.get('text', [])
        image_query = slide_data.get('imageQuery', '')
        
        is_title_slide = idx == 0
        is_end_slide = idx == total_slides - 1 and ("danke" in slide_title.lower() or "fragen" in slide_title.lower())
        has_text = len(text_points) > 0
        
        if is_title_slide or is_end_slide or not has_text:
            title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_para = title_frame.paragraphs[0]
            title_para.text = slide_title
            title_para.font.size = Pt(52)
            title_para.font.bold = True
            title_para.alignment = PP_ALIGN.CENTER
            
            if is_title_slide:
                subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4.8), Inches(9.333), Inches(1))
                sub_frame = subtitle_box.text_frame
                sub_para = sub_frame.paragraphs[0]
                sub_para.text = json_data.get('title', '')
                sub_para.font.size = Pt(26)
                sub_para.alignment = PP_ALIGN.CENTER
        else:
            title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.9))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = slide_title
            title_para.font.size = Pt(34)
            title_para.font.bold = True
            
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.2), Inches(5.5), Inches(0.05))
            line.fill.solid()
            line.line.fill.background()
            
            text_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.5))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            
            for i, point in enumerate(text_points):
                if i == 0:
                    para = text_frame.paragraphs[0]
                else:
                    para = text_frame.add_paragraph()
                para.text = f"•  {point}"
                para.font.size = Pt(20)
                para.space_before = Pt(6)
                para.space_after = Pt(12)
            
            if image_query:
                img_bytes, _ = fetch_pexels_image(image_query)
                if img_bytes:
                    try:
                        image_stream = io.BytesIO(img_bytes)
                        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.6), Inches(1.2), Inches(6.2), Inches(5.8))
                        frame.fill.solid()
                        frame.line.fill.background()
                        slide.shapes.add_picture(image_stream, Inches(6.7), Inches(1.3), Inches(6.0), Inches(5.6))
                    except Exception:
                        pass
    
    pptx_bytes = io.BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    return pptx_bytes.getvalue()

def render_presentation_preview(json_data):
    json_data = fix_presentation_data(json_data)
    st.markdown(f"### 📊 {json_data['title']}")
    
    for i, slide in enumerate(json_data['slides']):
        slide_title = slide.get('slideTitle', f'Folie {i+1}')
        text_items = slide.get('text', [])
        image_query = slide.get('imageQuery', '')
        
        if i == 0 or not text_items:
            st.markdown(f"""
                <div class="slide-preview" style="text-align:center;">
                    <h4 style="border:none; text-align:center;">Folie {i+1}: {slide_title}</h4>
                    <p style="color:#888; font-style:italic;">{'Titelfolie' if i == 0 else 'Abschlussfolie'}</p>
                </div>
            """, unsafe_allow_html=True)
        else:
            text_html = ''.join(f"<li>{p}</li>" for p in text_items)
            st.markdown(f"""
                <div class="slide-preview">
                    <h4>Folie {i+1}: {slide_title}</h4>
                    <div style="display:flex; gap:20px;">
                        <div style="flex:1;"><ul style="margin:0; padding-left:20px;">{text_html}</ul></div>
                        <div style="flex:1; text-align:center; padding:20px; background:#f5f5f5; border-radius:8px;">
                            <span style="color:#888;">🖼️ {image_query}</span>
                        </div>
                    </div>
                </div>
            """, unsafe_allow_html=True)

# ===== WELCOME SCREEN =====
if not st.session_state.welcome_done:
    show_confetti()
    
    st.markdown("<br><br><br>", unsafe_allow_html=True)
    
    if is_christmas_season():
        st.markdown("""
            <h1 class="welcome-title">🎄 Frohe Weihnachten, Sophia! 🎄</h1>
            <p class="welcome-subtitle">Ich hoffe das du dieses Geschenk magst!! </p>
        """, unsafe_allow_html=True)
        btn_text = "🎁 Geschenk öffnen"
    else:
        st.markdown("""
            <h1 class="welcome-title">👋 Willkommen, Sophia!</h1>
            <p class="welcome-subtitle">Dein persönlicher Lern-Assistent ist bereit!</p>
        """, unsafe_allow_html=True)
        btn_text = "🚀 Los geht's!"
    
    col1, col2, col3 = st.columns([1, 1, 1])
    with col2:
        if st.button(btn_text, use_container_width=True, type="primary"):
            st.session_state.welcome_done = True
            st.session_state.show_burst = True
            create_new_chat()
            do_rerun()
    
    st.stop()

# ===== CONFETTI BURST =====
if st.session_state.get('show_burst', False):
    show_confetti_burst()
    st.session_state.show_burst = False

# ===== SIDEBAR =====
with st.sidebar:
    st.markdown("## 💬 Chats")
    
    if st.button("➕ Neuer Chat", use_container_width=True, type="primary"):
        create_new_chat()
        do_rerun()
    
    st.markdown("---")
    
    sorted_chats = sorted(st.session_state.all_chats.items(), key=lambda x: x[0], reverse=True)
    
    for chat_id, chat_data in sorted_chats:
        is_active = chat_id == st.session_state.current_chat_id
        col1, col2 = st.columns([4, 1])
        
        with col1:
            title = chat_data.get("title", "Neuer Chat")
            icon = "📍" if is_active else "💭"
            if st.button(f"{icon} {title}", key=f"chat_{chat_id}", use_container_width=True):
                st.session_state.current_chat_id = chat_id
                st.session_state.current_pptx = chat_data.get("pptx_data")
                st.session_state.current_pptx_name = chat_data.get("pptx_name")
                do_rerun()
        
        with col2:
            if st.button("🗑️", key=f"del_{chat_id}"):
                del st.session_state.all_chats[chat_id]
                if chat_id == st.session_state.current_chat_id:
                    if st.session_state.all_chats:
                        st.session_state.current_chat_id = list(st.session_state.all_chats.keys())[0]
                    else:
                        create_new_chat()
                do_rerun()
    
    if not st.session_state.all_chats:
        create_new_chat()

# ===== MAIN APP =====
if st.session_state.current_chat_id is None or st.session_state.current_chat_id not in st.session_state.all_chats:
    create_new_chat()

current_chat = st.session_state.all_chats[st.session_state.current_chat_id]

st.markdown("""
    <div class="header-container">
        <h1>🎓 Deine Lern-Assistent</h1>
        <p>Hausaufgabenhilfe & PowerPoint-Präsentationen!</p>
    </div>
""", unsafe_allow_html=True)

messages = current_chat.get("messages", [])

if not messages:
    st.markdown("""
        <div class="chat-message ai">
            <span class="avatar">🤖</span>
            <div class="message-content">
                Hallo Sophia! 👋<br><br>
                Ich bin eine KI, die von <b>Shiva</b> als Weihnachtsgeschenk für dich entwickelt wurde! 🎁<br><br>
                <b>Ich kann:</b><br>
                📚 Bei Hausaufgaben helfen<br>
                📊 PowerPoint-Präsentationen erstellen<br><br>
                Sag z.B.: <i>"Erstelle eine Präsentation über Vulkane"</i>
            </div>
        </div>
    """, unsafe_allow_html=True)

for msg in messages:
    if msg["role"] == "user":
        st.markdown(f"""
            <div class="chat-message user">
                <span class="avatar">👩</span>
                <div class="message-content">{msg["content"]}</div>
            </div>
        """, unsafe_allow_html=True)
    else:
        if msg.get("is_presentation"):
            st.markdown("""
                <div class="chat-message ai">
                    <span class="avatar">🤖</span>
                    <div class="message-content">✅ Deine Präsentation ist fertig! Hier ist die Vorschau:</div>
                </div>
            """, unsafe_allow_html=True)
            render_presentation_preview(msg["presentation_data"])
        else:
            st.markdown(f"""
                <div class="chat-message ai">
                    <span class="avatar">🤖</span>
                    <div class="message-content">{msg["content"]}</div>
                </div>
            """, unsafe_allow_html=True)

# ===== DOWNLOAD BUTTON =====
if st.session_state.current_pptx is not None and st.session_state.current_pptx_name is not None:
    st.markdown("""
        <div class="download-section">
            <h3 style="color:white; margin:0;">📥 Deine PowerPoint ist bereit!</h3>
        </div>
    """, unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        st.download_button(
            label="⬇️ PowerPoint herunterladen (.pptx)",
            data=st.session_state.current_pptx,
            file_name=st.session_state.current_pptx_name,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
            type="primary"
        )

# ===== SEND MESSAGE FUNCTION =====
def send_message():
    user_input = st.session_state.get(f"user_input_{st.session_state.input_key}", "")
    
    if user_input.strip():
        current_chat["messages"].append({"role": "user", "content": user_input})
        
        if len(current_chat["messages"]) == 1:
            current_chat["title"] = get_chat_title(current_chat["messages"])
        
        st.session_state.current_pptx = None
        st.session_state.current_pptx_name = None
        current_chat["pptx_data"] = None
        current_chat["pptx_name"] = None
        
        response = call_groq_api(current_chat["messages"])
        
        try:
            json_start = response.find('{')
            json_end = response.rfind('}') + 1
            
            if json_start != -1 and json_end > json_start:
                json_str = response[json_start:json_end]
                json_data = json.loads(json_str)
                
                if "title" in json_data and "slides" in json_data:
                    current_chat["messages"].append({
                        "role": "assistant",
                        "content": "",
                        "is_presentation": True,
                        "presentation_data": json_data
                    })
                    
                    pptx_bytes = create_powerpoint(json_data)
                    
                    if pptx_bytes:
                        safe_title = json_data['title']
                        for char in ['/', '\\', ':', '*', '?', '"', '<', '>', '|']:
                            safe_title = safe_title.replace(char, '_')
                        filename = f"{safe_title[:50]}.pptx"
                        
                        st.session_state.current_pptx = pptx_bytes
                        st.session_state.current_pptx_name = filename
                        current_chat["pptx_data"] = pptx_bytes
                        current_chat["pptx_name"] = filename
                else:
                    current_chat["messages"].append({"role": "assistant", "content": response})
            else:
                current_chat["messages"].append({"role": "assistant", "content": response})
                
        except json.JSONDecodeError:
            current_chat["messages"].append({"role": "assistant", "content": response})
        
        st.session_state.input_key += 1

# ===== INPUT =====
st.markdown("<br>", unsafe_allow_html=True)
col1, col2 = st.columns([5, 1])

with col1:
    st.text_input(
        "Nachricht",
        placeholder="Schreibe hier... (Enter zum Senden)",
        label_visibility="collapsed",
        key=f"user_input_{st.session_state.input_key}",
        on_change=send_message
    )

with col2:
    if st.button("📤", use_container_width=True, type="primary"):
        send_message()
        do_rerun()

st.markdown("""
    <br><br>
    <div style="text-align:center; color:rgba(255,255,255,0.4); font-size:0.8rem;">
        von Shiva:)    </div>
""", unsafe_allow_html=True)





